#!/usr/bin/env python3
"""Generate CarManage customer/owner flow PowerPoint with diagrams."""

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Install: pip install python-pptx")
    raise

OUT = Path(__file__).resolve().parent.parent / "CarManage_Flow_Diagrams.pptx"

BLUE = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_title(slide, title: str, subtitle: str = ""):
    if slide.shapes.title:
        slide.shapes.title.text = title
        for p in slide.shapes.title.text_frame.paragraphs:
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = BLUE
    if subtitle and len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        ph.text = subtitle
        for p in ph.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY


def add_bullets(slide, items: list[str], left=0.8, top=1.8, width=8.5, height=5):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(8)


def add_flow_box(slide, text: str, x, y, w=1.6, h=0.55, fill=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BLUE
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, x1, y1, x2, y2):
    slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)).line.color.rgb = GRAY


def slide_title(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_title(slide, title, subtitle)
    return slide


def slide_content(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_title(slide, title)
    add_bullets(slide, bullets)
    return slide


def build_customer_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Customer Flow"
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.size = Pt(28)
        p.font.color.rgb = BLUE

    steps = [
        ("Register +\nKYC + License", 0.5, 1.2),
        ("Search &\nBook", 2.3, 1.2),
        ("Negotiate\n(Chat)", 4.1, 1.2),
        ("Accept\nPrice", 5.9, 1.2),
        ("Pay 75%\nDeposit", 7.7, 1.2),
        ("Pickup\nCheck-in", 0.5, 2.8),
        ("On\nTrip", 2.3, 2.8),
        ("Return\nVehicle", 4.1, 2.8),
        ("Pay Final\nBalance", 5.9, 2.8),
        ("Review", 7.7, 2.8),
    ]
    for text, x, y in steps:
        add_flow_box(slide, text, x, y)
    # row 1 arrows
    for i in range(4):
        add_arrow(slide, 0.5 + 1.8 * (i + 1) - 0.15, 1.45, 0.5 + 1.8 * (i + 1) + 0.05, 1.45)
    add_arrow(slide, 8.5, 1.75, 1.0, 2.55)
    for i in range(3):
        add_arrow(slide, 0.5 + 1.8 * (i + 1) - 0.15, 3.05, 0.5 + 1.8 * (i + 1) + 0.05, 3.05)


def build_owner_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Owner Flow"
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.size = Pt(28)
        p.font.color.rgb = BLUE

    steps = [
        ("Register +\nKYC", 0.5, 1.2),
        ("Fleet:\nAdd Car", 2.3, 1.2),
        ("Booking\nRequest", 4.1, 1.2),
        ("Set Price &\nConfirm", 5.9, 1.2),
        ("Vehicle\nHandover", 7.7, 1.2),
        ("Post-Trip\nCharges", 2.3, 2.8),
        ("Accept\nReturn", 4.1, 2.8),
        ("Payout\nView", 5.9, 2.8),
        ("Review", 7.7, 2.8),
    ]
    for text, x, y in steps:
        add_flow_box(slide, text, x, y)
    for i in range(4):
        add_arrow(slide, 0.5 + 1.8 * (i + 1) - 0.15, 1.45, 0.5 + 1.8 * (i + 1) + 0.05, 1.45)
    add_arrow(slide, 8.5, 1.75, 2.8, 2.55)
    for i in range(2):
        add_arrow(slide, 2.3 + 1.8 * (i + 1) - 0.15, 3.05, 2.3 + 1.8 * (i + 1) + 0.05, 3.05)


def build_combined_sequence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Combined Booking Lifecycle"
    for p in slide.shapes.title.text_frame.paragraphs:
        p.font.size = Pt(26)
        p.font.color.rgb = BLUE

    cols = [
        ("Customer", 0.4, ["Inquiry", "Accept price", "75% deposit", "Pickup/Return", "Final pay", "Review"]),
        ("Platform", 3.5, ["Overlap check", "Fees & GST", "Settlement", "Payment gate", "Complete"]),
        ("Owner", 6.6, ["Set price", "Confirm", "Handover", "Post-trip charges", "Accept return", "Review"]),
    ]
    for label, x, items in cols:
        add_flow_box(slide, label, x, 1.0, w=2.6, h=0.45, fill=BLUE)
        for i, item in enumerate(items):
            add_flow_box(slide, item, x, 1.65 + i * 0.65, w=2.6, h=0.5, fill=ACCENT if label != "Platform" else GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs, "CarManage (CM-P2P)", "Customer & Owner Flow Diagrams\nPeer-to-peer car rental platform")

    slide_content(
        prs,
        "Overview",
        [
            "Roles: CUSTOMER and OWNER only (no admin console)",
            "Stack: Go/Gin API + Next.js frontend + PostgreSQL",
            "Payments: 75% deposit, then post-trip charges, then final balance",
            "Pricing: Negotiated per-day rate × inclusive trip days + platform fees/GST",
            "Docs: docs/CUSTOMER_FLOW.md and docs/OWNER_FLOW.md",
        ],
    )

    slide_content(
        prs,
        "Booking Status States",
        [
            "PENDING → customer creates inquiry",
            "NEGOTIATING → chat / price discussion",
            "CONFIRMED → customer accepted price + owner confirmed",
            "COMPLETED → owner accepted return after full payment",
            "CANCELLED → withdraw, cancel, or unpaid cancel",
        ],
    )

    slide_content(
        prs,
        "Payment States (on CONFIRMED)",
        [
            "UNPAID → customer pays 75% deposit",
            "DEPOSIT_PAID → trip handover begins",
            "FINAL_DUE → owner submitted post-trip charges",
            "PAID → customer paid final balance; owner can accept return",
        ],
    )

    build_customer_flow(prs)
    build_owner_flow(prs)
    build_combined_sequence(prs)

    slide_content(
        prs,
        "Calculations — Customer Total",
        [
            "Trip days = inclusive UTC calendar days (min 1)",
            "Agreed base = FinalBookingPrice × trip days",
            "Customer commission = agreed base × 2% (configurable)",
            "Customer GST = (agreed base + commission) × 18%",
            "Customer total = base + commission + GST",
            "Deposit = 75% of customer total | Final = total + post-trip − deposit",
        ],
    )

    slide_content(
        prs,
        "Calculations — Owner Net",
        [
            "Owner commission = agreed base × 1.5% (configurable)",
            "Owner GST = agreed base × 18%",
            "Owner net = agreed base − owner commission − owner GST",
            "Projected payout = owner net + post-trip charges (until fully paid)",
            "Listing per-hour/per-km prices are hints only — not used in settlement",
        ],
    )

    slide_content(
        prs,
        "Example (₹1000/day × 2 days)",
        [
            "Agreed base: ₹2,000.00",
            "Customer total: ₹2,407.20 (2% fee + 18% GST on base+fee)",
            "Deposit (75%): ₹1,805.40",
            "Owner net: ₹1,610.00 (1.5% fee + 18% GST on base)",
            "If post-trip charges ₹500 → final due: ₹1,101.80",
        ],
    )

    slide_content(
        prs,
        "Key Validations — Customer",
        [
            "KYC verified + driving license before booking",
            "acknowledged_deposit_terms on create",
            "No overlapping bookings; cannot book own car",
            "Accept price only after owner sets final_booking_price",
            "Final pay only when status FINAL_DUE (post-trip charges submitted)",
            "Handover: odometer 1–2M km; return ≥ pickup",
        ],
    )

    slide_content(
        prs,
        "Key Validations — Owner",
        [
            "KYC verified for fleet and booking management routes",
            "Car: fuel/transmission/year/seats/airbag rules",
            "Cannot edit car booked for current UTC day",
            "Confirm only after customer accepted price",
            "Post-trip: max 30 lines, total ≤ ₹20 lakh, after customer return",
            "Accept return only after payment PAID",
        ],
    )

    slide_content(
        prs,
        "Trip Handover Stages",
        [
            "1. Awaiting owner handover (after deposit)",
            "2. Awaiting customer pickup",
            "3. On trip",
            "4. Awaiting customer return",
            "5. Awaiting post-trip charges (owner)",
            "6. Awaiting final payment (customer)",
            "7. Awaiting owner return acceptance",
            "8. Trip complete",
        ],
    )

    prs.save(OUT)
    print(f"Created: {OUT}")


if __name__ == "__main__":
    main()
